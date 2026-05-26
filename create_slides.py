from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ---- Color Palette ----
BLUE_DARK   = RGBColor(0x1E, 0x3A, 0x5F)   # navy
BLUE_MID    = RGBColor(0x23, 0x6E, 0xC4)   # brand blue
BLUE_LIGHT  = RGBColor(0xDB, 0xEA, 0xFE)   # very light blue
EMERALD     = RGBColor(0x05, 0x96, 0x69)   # emerald
EMERALD_LT  = RGBColor(0xD1, 0xFA, 0xE5)   # light emerald
AMBER       = RGBColor(0xD9, 0x77, 0x06)
AMBER_LT    = RGBColor(0xFE, 0xF3, 0xC7)
PURPLE      = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_LT   = RGBColor(0xED, 0xE9, 0xFE)
SLATE_DARK  = RGBColor(0x1E, 0x29, 0x3B)
SLATE_MID   = RGBColor(0x47, 0x55, 0x69)
SLATE_LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT  = RGBColor(0xE2, 0xE8, 0xF0)

blank_layout = prs.slide_layouts[6]   # completely blank


# ════════════════════════════════════════════════════════
# Helpers
# ════════════════════════════════════════════════════════

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
             wrap=True, font="Sarabun"):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txb


def add_bullet_box(slide, items, l, t, w, h,
                   size=16, color=None, bullet="▸", line_spacing=1.15):
    from pptx.oxml.ns import qn
    from lxml import etree
    import copy

    color = color or SLATE_DARK
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if isinstance(item, tuple):
            prefix, rest = item
        else:
            prefix, rest = bullet + " ", item

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.alignment = PP_ALIGN.LEFT
        # space before each bullet
        p.space_before = Pt(4)

        run = p.add_run()
        run.text = prefix + rest
        run.font.name = "Sarabun"
        run.font.size = Pt(size)
        run.font.color.rgb = color

    return txb


def section_header(slide, label, color=BLUE_MID):
    rect = add_rect(slide, 0, 0, 13.33, 0.55, fill=color)
    add_text(slide, label, 0.4, 0.08, 12, 0.45,
             size=14, bold=True, color=WHITE, align=PP_ALIGN.LEFT)


def slide_title(slide, title, subtitle=None,
                title_color=SLATE_DARK, sub_color=SLATE_MID,
                title_size=32, sub_size=18):
    add_text(slide, title, 0.6, 0.65, 12, 0.9,
             size=title_size, bold=True, color=title_color)
    if subtitle:
        add_text(slide, subtitle, 0.6, 1.45, 12, 0.5,
                 size=sub_size, bold=False, color=sub_color)


def add_icon_card(slide, icon, label, desc, l, t, w=2.8, h=1.6,
                  bg=BLUE_LIGHT, ic_color=BLUE_MID):
    add_rect(slide, l, t, w, h, fill=bg)
    add_text(slide, icon, l+0.15, t+0.1, 0.55, 0.55, size=26, align=PP_ALIGN.CENTER)
    add_text(slide, label, l+0.75, t+0.1, w-0.85, 0.4,
             size=14, bold=True, color=SLATE_DARK)
    add_text(slide, desc,  l+0.75, t+0.48, w-0.85, 0.9,
             size=11, color=SLATE_MID, wrap=True)


# ════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)

# full gradient-like background (two rects)
add_rect(slide, 0, 0, 13.33, 7.5, fill=RGBColor(0xF0, 0xF7, 0xFF))
add_rect(slide, 0, 0, 5.2, 7.5, fill=BLUE_DARK)

# left white accent stripe
add_rect(slide, 4.9, 0, 0.3, 7.5, fill=BLUE_MID)

# left side content
add_text(slide, "Prompt-Berk", 0.5, 2.2, 4.2, 1.2,
         size=48, bold=True, color=WHITE)
add_text(slide, "ระบบบริหารจัดการ\nงบประมาณและเบิกจ่าย", 0.5, 3.4, 4.2, 1.4,
         size=20, color=BLUE_LIGHT, wrap=True)
add_text(slide, "สถาบันรับรองคุณภาพสถานพยาบาล (สรพ.)  |  2026",
         0.5, 6.8, 4.5, 0.45, size=11, color=BLUE_LIGHT)

# right side summary boxes
cards = [
    ("💸", "เบิกจ่าย",   "สร้างรายการเบิกได้ทันที"),
    ("📋", "ติดตาม",     "ดูประวัติและสถานะ"),
    ("⚙️", "Admin",      "จัดการหมวดหมู่งบ"),
    ("🐳", "Docker",     "Deploy บนเซิร์ฟเวอร์"),
]
col1_x = 5.6
col2_x = 9.3
row1_y = 1.2
row2_y = 4.0

positions = [(col1_x, row1_y), (col2_x, row1_y),
             (col1_x, row2_y), (col2_x, row2_y)]

for (icon, lbl, desc), (lx, ly) in zip(cards, positions):
    add_rect(slide, lx, ly, 3.3, 2.4, fill=WHITE)
    add_text(slide, icon, lx+0.25, ly+0.3, 0.7, 0.7, size=30, align=PP_ALIGN.CENTER)
    add_text(slide, lbl,  lx+1.1,  ly+0.3, 2.0, 0.5, size=18, bold=True, color=BLUE_DARK)
    add_text(slide, desc, lx+1.1,  ly+0.8, 2.0, 0.9, size=13, color=SLATE_MID, wrap=True)


# ════════════════════════════════════════════════════════
# SLIDE 2 — ปัญหาเดิม
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=SLATE_LIGHT)
section_header(slide, "01  |  ที่มาของปัญหา", color=SLATE_DARK)

slide_title(slide, "ก่อนจะมี Prompt-Berk...",
            subtitle="การเบิกจ่ายงบประมาณใช้กระบวนการแบบเดิมที่ใช้เวลาและเกิดข้อผิดพลาดได้ง่าย")

problems = [
    ("📄  ", "กรอกแบบฟอร์มกระดาษหรือ Excel ด้วยมือ — ทำซ้ำทุกครั้ง"),
    ("📧  ", "ส่งเอกสารผ่านอีเมลหลายรอบ — ติดตามสถานะยาก"),
    ("🔢  ", "คำนวณยอดเงินเองทุกครั้ง — เสี่ยงต่อความผิดพลาด"),
    ("🗂️  ", "ข้อมูลกระจัดกระจายในหลายไฟล์ — หาย้อนหลังยาก"),
    ("⏱️  ", "ใช้เวลานานในการรวบรวมข้อมูลเพื่อรายงาน"),
]

add_bullet_box(slide, problems, 1.0, 2.1, 11.0, 4.5,
               size=17, color=SLATE_DARK, bullet="")

# bottom bar
add_rect(slide, 0, 6.8, 13.33, 0.7, fill=BLUE_MID)
add_text(slide, "Prompt-Berk ถูกสร้างขึ้นเพื่อแก้ปัญหาเหล่านี้ทั้งหมด",
         1.0, 6.88, 11.0, 0.5, size=16, bold=True, color=WHITE)


# ════════════════════════════════════════════════════════
# SLIDE 3 — Prompt-Berk คืออะไร
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=RGBColor(0xF0, 0xF7, 0xFF))
section_header(slide, "02  |  ภาพรวมระบบ", color=BLUE_MID)

slide_title(slide, "Prompt-Berk คืออะไร?")

points = [
    ("🌐  ", "ระบบเว็บสำหรับบริหารจัดการงบประมาณและเบิกจ่ายของ สรพ."),
    ("🖥️  ", "ใช้งานผ่านเบราว์เซอร์ทุกอุปกรณ์ — ไม่ต้องติดตั้งโปรแกรม"),
    ("🗄️  ", "ข้อมูลเก็บส่วนกลางบนเซิร์ฟเวอร์ — ทุกคนเห็นข้อมูลร่วมกัน"),
    ("🏢  ", "Deploy บนเซิร์ฟเวอร์ภายในองค์กร (Intranet) ด้วย Docker"),
    ("⚡  ", "คำนวณยอดเงินอัตโนมัติ ลดความผิดพลาด"),
]
add_bullet_box(slide, points, 1.0, 2.05, 7.5, 4.5,
               size=17, color=SLATE_DARK, bullet="")

# right info box
add_rect(slide, 9.0, 1.8, 3.8, 5.0, fill=BLUE_DARK)
add_text(slide, "Tech Stack", 9.3, 2.0, 3.2, 0.5,
         size=14, bold=True, color=BLUE_LIGHT)

stack = [
    "Frontend",  "HTML + Tailwind CSS",
    "Backend",   "Node.js + Express",
    "Database",  "SQLite",
    "Server",    "nginx (Docker)",
]
for i in range(0, len(stack), 2):
    y = 2.6 + (i//2) * 0.85
    add_text(slide, stack[i],   9.3, y,      1.4, 0.4, size=11, color=BLUE_LIGHT, bold=True)
    add_text(slide, stack[i+1], 10.7, y,     1.9, 0.4, size=11, color=WHITE)


# ════════════════════════════════════════════════════════
# SLIDE 4 — ผู้ใช้งาน
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=SLATE_LIGHT)
section_header(slide, "03  |  ผู้ใช้งาน", color=EMERALD)

slide_title(slide, "ผู้ใช้งานระบบมีใครบ้าง?",
            title_color=SLATE_DARK)

user_cards = [
    ("👩‍💼", "เจ้าหน้าที่ผู้เบิก",
     ["สร้างรายการเบิกจ่าย", "ติดตามสถานะรายการ", "พิมพ์เอกสาร PDF"],
     EMERALD_LT, EMERALD),
    ("👨‍💻", "Admin / ผู้ประสานงาน",
     ["จัดการหมวดหมู่งบประมาณ", "ตั้งค่าสูตรคำนวณ", "ดูข้อมูลรวม"],
     BLUE_LIGHT, BLUE_MID),
    ("🔧", "ทีม IT",
     ["ติดตั้ง Docker บนเซิร์ฟเวอร์", "ดูแลระบบและสำรองข้อมูล", "จัดการ Network/Firewall"],
     AMBER_LT, AMBER),
]

for i, (icon, title, bullets, bg, accent) in enumerate(user_cards):
    lx = 0.5 + i * 4.25
    add_rect(slide, lx, 2.0, 3.9, 4.8, fill=bg)
    add_rect(slide, lx, 2.0, 3.9, 0.55, fill=accent)
    add_text(slide, icon,  lx+0.2, 2.0, 0.55, 0.55,  size=22, color=WHITE)
    add_text(slide, title, lx+0.8, 2.07, 3.0, 0.45,  size=15, bold=True, color=WHITE)
    add_bullet_box(slide, bullets, lx+0.3, 2.75, 3.4, 3.5,
                   size=14, color=SLATE_DARK, bullet="•")


# ════════════════════════════════════════════════════════
# SLIDE 5 — สร้างรายการเบิก
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=EMERALD_LT)
section_header(slide, "04  |  ฟีเจอร์ที่ 1", color=EMERALD)

slide_title(slide, "💸  สร้างรายการเบิกใหม่",
            subtitle="ขออนุมัติเบิกงบประมาณ เลือกหมวดหมู่และกิจกรรม แล้วระบบคำนวณยอดให้อัตโนมัติ",
            title_color=SLATE_DARK, sub_color=SLATE_MID)

steps = [
    ("1", "เลือกหมวดหมู่งบประมาณ",   "เลือกจากรายการที่ Admin ตั้งไว้"),
    ("2", "เลือกกิจกรรมและรายการ",    "ระบบโหลดสูตรคำนวณให้อัตโนมัติ"),
    ("3", "กรอกข้อมูลการเดินทาง",     "ดึงข้อมูลจากโปรไฟล์ส่วนตัว (Autofill)"),
    ("4", "คำนวณระยะทาง (Longdo Map)", "เลือกจุดเริ่มต้น-ปลายทางบนแผนที่"),
    ("5", "บันทึกและส่งรายการ",        "ข้อมูลถูกเก็บบนเซิร์ฟเวอร์ส่วนกลาง"),
]

for i, (num, title, desc) in enumerate(steps):
    lx = 0.6
    ly = 2.1 + i * 0.95
    add_rect(slide, lx, ly, 0.55, 0.55, fill=EMERALD)
    add_text(slide, num, lx, ly, 0.55, 0.55,
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, lx+0.7, ly+0.03, 4.5, 0.35,
             size=15, bold=True, color=SLATE_DARK)
    add_text(slide, desc,  lx+0.7, ly+0.35, 4.5, 0.35,
             size=12, color=SLATE_MID)

# right callout
add_rect(slide, 7.0, 2.0, 5.8, 4.8, fill=WHITE)
add_rect(slide, 7.0, 2.0, 5.8, 0.5, fill=EMERALD)
add_text(slide, "✨  จุดเด่น", 7.2, 2.05, 5.2, 0.4,
         size=14, bold=True, color=WHITE)
highlights = [
    "• คำนวณค่าเดินทางอัตโนมัติ (กม. × อัตรา)",
    "• Autofill จากโปรไฟล์ส่วนตัว",
    "• ระบุตำแหน่งผ่านแผนที่ Longdo Map",
    "• บันทึกทันที ไม่ต้องรอส่งอีเมล",
]
add_bullet_box(slide, highlights, 7.2, 2.65, 5.3, 4.0,
               size=14, color=SLATE_DARK, bullet="")


# ════════════════════════════════════════════════════════
# SLIDE 6 — ประวัติการเบิกจ่าย
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=PURPLE_LT)
section_header(slide, "05  |  ฟีเจอร์ที่ 2", color=PURPLE)

slide_title(slide, "📋  ประวัติการเบิกจ่าย",
            subtitle="ตรวจสอบสถานะรายการ อัปเดตข้อมูล และส่งออกเอกสาร PDF",
            title_color=SLATE_DARK, sub_color=SLATE_MID)

features = [
    ("🔍", "ดูรายการเบิกทั้งหมด",    "แสดงรายการที่ทุกคนส่งบนเซิร์ฟเวอร์"),
    ("✏️", "แก้ไขรายการ",            "อัปเดตข้อมูลได้หลังบันทึก"),
    ("🖨️", "พิมพ์ PDF",              "สร้างเอกสารสรุปสำหรับใช้เบิกจ่ายจริง"),
    ("🗑️", "ลบรายการ",              "จัดการรายการที่ไม่ต้องการ"),
]

for i, (icon, title, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    lx = 0.8 + col * 6.1
    ly = 2.2 + row * 2.2
    add_rect(slide, lx, ly, 5.7, 1.8, fill=WHITE)
    add_text(slide, icon,  lx+0.25, ly+0.5,  0.6, 0.6, size=24, align=PP_ALIGN.CENTER)
    add_text(slide, title, lx+1.0,  ly+0.2,  4.3, 0.5, size=16, bold=True, color=SLATE_DARK)
    add_text(slide, desc,  lx+1.0,  ly+0.75, 4.3, 0.7, size=13, color=SLATE_MID, wrap=True)


# ════════════════════════════════════════════════════════
# SLIDE 7 — โปรไฟล์ส่วนตัว
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=RGBColor(0xEE, 0xF2, 0xFF))
section_header(slide, "06  |  ฟีเจอร์ที่ 3", color=RGBColor(0x43, 0x38, 0xCA))

slide_title(slide, "👤  โปรไฟล์ส่วนตัว",
            subtitle="ตั้งค่าข้อมูลผู้เดินทางครั้งเดียว ระบบจะ Autofill ทุกครั้งที่สร้างรายการเบิก",
            title_color=SLATE_DARK, sub_color=SLATE_MID)

# two column layout
left_items = [
    "ชื่อ-นามสกุล และตำแหน่ง",
    "สังกัดหน่วยงาน",
    "ที่อยู่ปัจจุบัน (สำหรับคำนวณระยะทาง)",
    "อัตราค่าเดินทางส่วนตัว",
]
right_note = "ข้อมูลโปรไฟล์เก็บใน LocalStorage\nของเบราว์เซอร์บนเครื่องตัวเอง\n(ไม่ได้อัปโหลดขึ้นเซิร์ฟเวอร์)"

add_text(slide, "ข้อมูลที่ตั้งค่าได้:", 0.8, 2.1, 7.0, 0.45,
         size=16, bold=True, color=SLATE_DARK)
add_bullet_box(slide, left_items, 0.8, 2.6, 7.0, 3.5,
               size=16, color=SLATE_DARK)

add_rect(slide, 8.8, 2.0, 4.0, 3.5, fill=AMBER_LT)
add_rect(slide, 8.8, 2.0, 4.0, 0.45, fill=AMBER)
add_text(slide, "📌  หมายเหตุ", 9.0, 2.05, 3.6, 0.4,
         size=13, bold=True, color=WHITE)
add_text(slide, right_note, 9.0, 2.6, 3.6, 2.7,
         size=13, color=SLATE_DARK, wrap=True)


# ════════════════════════════════════════════════════════
# SLIDE 8 — Admin
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=SLATE_LIGHT)
section_header(slide, "07  |  ฟีเจอร์ที่ 4", color=SLATE_DARK)

slide_title(slide, "⚙️  ตั้งค่าระบบ (Admin)",
            subtitle="จัดการหมวดหมู่งบประมาณและสูตรคำนวณ — ข้อมูลนี้ใช้ร่วมกันทั้งระบบ",
            title_color=SLATE_DARK, sub_color=SLATE_MID)

admin_items = [
    ("➕", "เพิ่มหมวดหมู่งบประมาณ",  "กำหนดชื่อหมวดและรายการย่อย"),
    ("✏️", "แก้ไขสูตรคำนวณ",         "ตั้งอัตราค่าใช้จ่ายต่อหน่วย"),
    ("🗑️", "ลบหมวดหมู่",             "ลบรายการที่ไม่ต้องการออก"),
    ("📦", "โหลดข้อมูลตัวอย่าง",     "Demo data สำหรับทดสอบระบบ"),
]

for i, (icon, title, desc) in enumerate(admin_items):
    lx = 0.7 + (i % 2) * 6.2
    ly = 2.2 + (i // 2) * 2.2
    add_rect(slide, lx, ly, 5.8, 1.8, fill=WHITE)
    add_rect(slide, lx, ly, 0.6, 1.8, fill=SLATE_DARK)
    add_text(slide, icon, lx+0.05, ly+0.55, 0.5, 0.6, size=20, align=PP_ALIGN.CENTER, color=WHITE)
    add_text(slide, title, lx+0.8, ly+0.2,  4.7, 0.5, size=15, bold=True, color=SLATE_DARK)
    add_text(slide, desc,  lx+0.8, ly+0.75, 4.7, 0.7, size=13, color=SLATE_MID, wrap=True)


# ════════════════════════════════════════════════════════
# SLIDE 9 — ข้อมูลเก็บที่ไหน
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=RGBColor(0xF0, 0xF7, 0xFF))
section_header(slide, "08  |  การจัดเก็บข้อมูล", color=BLUE_MID)

slide_title(slide, "ข้อมูลถูกเก็บที่ไหน?",
            title_color=SLATE_DARK)

# Table header
add_rect(slide, 0.6, 2.05, 12.0, 0.55, fill=BLUE_DARK)
cols = ["ข้อมูล", "เก็บที่", "ใครเห็น"]
col_widths = [4.5, 4.0, 3.5]
col_x = [0.6, 5.1, 9.1]

for label, cx, cw in zip(cols, col_x, col_widths):
    add_text(slide, label, cx+0.15, 2.1, cw-0.2, 0.45,
             size=14, bold=True, color=WHITE)

rows = [
    ("หมวดหมู่งบประมาณ (Admin)", "SQLite บนเซิร์ฟเวอร์", "ทุกคนในระบบ"),
    ("รายการเบิกจ่าย",           "SQLite บนเซิร์ฟเวอร์", "ทุกคนในระบบ"),
    ("โปรไฟล์ส่วนตัว",          "LocalStorage (Browser)", "เฉพาะเครื่องตัวเอง"),
]

row_colors = [WHITE, SLATE_LIGHT, WHITE]
for ri, (row, bg) in enumerate(zip(rows, row_colors)):
    ry = 2.6 + ri * 1.35
    add_rect(slide, 0.6, ry, 12.0, 1.3, fill=bg)
    for val, cx, cw in zip(row, col_x, col_widths):
        color = EMERALD if "SQLite" in val else (AMBER if "Local" in val else SLATE_DARK)
        add_text(slide, val, cx+0.15, ry+0.38, cw-0.2, 0.6,
                 size=14, color=color, bold=("SQLite" in val or "Local" in val))


# ════════════════════════════════════════════════════════
# SLIDE 10 — Architecture
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=SLATE_DARK)
section_header(slide, "09  |  Architecture (IT)", color=BLUE_MID)

slide_title(slide, "Architecture ของระบบ",
            subtitle="ระบบประกอบด้วย 2 container ใน Docker ทำงานร่วมกัน",
            title_color=WHITE, sub_color=BLUE_LIGHT)

# Architecture diagram boxes
boxes = [
    (0.5, 3.2, 2.4, 2.0, "🌐", "ผู้ใช้งาน", "เบราว์เซอร์", BLUE_LIGHT, BLUE_MID),
    (3.5, 3.2, 2.8, 2.0, "nginx", "Frontend", "port 8080", RGBColor(0xBB, 0xF7, 0xD0), EMERALD),
    (7.0, 3.2, 2.8, 2.0, "Node.js", "Backend", "port 3000", AMBER_LT, AMBER),
    (10.4, 3.2, 2.4, 2.0, "🗄️", "SQLite", "Docker Volume", RGBColor(0xFE, 0xCA, 0xCA), RGBColor(0xDC, 0x26, 0x26)),
]

for lx, ly, lw, lh, icon, title, sub, bg, accent in boxes:
    add_rect(slide, lx, ly, lw, lh, fill=bg)
    add_rect(slide, lx, ly, lw, 0.45, fill=accent)
    add_text(slide, icon,  lx+0.1,     ly+0.05, lw-0.2, 0.4,  size=13, bold=True, color=WHITE)
    add_text(slide, title, lx+0.15,    ly+0.65, lw-0.2, 0.55, size=16, bold=True, color=SLATE_DARK)
    add_text(slide, sub,   lx+0.15,    ly+1.2,  lw-0.2, 0.5,  size=12, color=SLATE_MID)

# Arrows between boxes
arrow_positions = [
    (2.95, 4.22), (6.35, 4.22), (9.85, 4.22),
]
for ax, ay in arrow_positions:
    add_rect(slide, ax, ay+0.05, 0.55, 0.2, fill=BLUE_LIGHT)
    add_text(slide, "→", ax+0.05, ay-0.05, 0.5, 0.4, size=16, bold=True, color=WHITE)

# Labels under arrows
arrow_labels = [
    (2.9, 4.7, "HTTP :8080"),
    (6.25, 4.7, "Proxy API"),
    (9.75, 4.7, "SQL"),
]
for lx, ly, lbl in arrow_labels:
    add_text(slide, lbl, lx, ly, 0.8, 0.35, size=10, color=BLUE_LIGHT)

# Docker border note
add_rect(slide, 3.2, 2.85, 9.8, 3.0, fill=None, line=BLUE_MID, line_w=1.5)
add_text(slide, "🐳  Docker Environment", 3.35, 2.9, 3.5, 0.35,
         size=11, color=BLUE_LIGHT)


# ════════════════════════════════════════════════════════
# SLIDE 11 — วิธีติดตั้ง
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=SLATE_DARK)
section_header(slide, "10  |  การติดตั้ง (IT)", color=EMERALD)

slide_title(slide, "วิธีติดตั้งบนเซิร์ฟเวอร์",
            subtitle="ต้องการแค่ Docker + Docker Compose — ใช้ได้ทันทีใน 3 ขั้นตอน",
            title_color=WHITE, sub_color=BLUE_LIGHT)

install_steps = [
    ("1", "Clone หรือ copy โค้ดลงเซิร์ฟเวอร์",
     "git clone / copy ไฟล์โปรเจกต์ทั้งหมดไปยังเซิร์ฟเวอร์"),
    ("2", "รัน Docker Compose",
     "docker compose up -d --build"),
    ("3", "เปิดเบราว์เซอร์",
     "http://<IP-เซิร์ฟเวอร์>:8080/"),
]

for i, (num, title, cmd) in enumerate(install_steps):
    ly = 2.2 + i * 1.55
    add_rect(slide, 0.6, ly, 0.65, 0.65, fill=EMERALD)
    add_text(slide, num, 0.6, ly, 0.65, 0.65,
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, 1.45, ly+0.03, 5.5, 0.45,
             size=15, bold=True, color=WHITE)
    add_rect(slide, 1.45, ly+0.5, 11.0, 0.6, fill=RGBColor(0x0F, 0x17, 0x2A))
    add_text(slide, cmd, 1.65, ly+0.55, 10.6, 0.45,
             size=13, color=EMERALD, font="Courier New")

# right tip
add_rect(slide, 7.5, 2.1, 5.4, 4.5, fill=RGBColor(0x14, 0x23, 0x3A))
add_text(slide, "💡  ข้อมูลไม่หายเมื่อ Rebuild", 7.7, 2.2, 5.0, 0.5,
         size=13, bold=True, color=EMERALD)
add_text(slide,
         "ฐานข้อมูล SQLite เก็บใน\nDocker Volume ชื่อ db_data\n\n"
         "แม้จะรัน rebuild ข้อมูลทั้งหมด\nยังคงอยู่ครบถ้วน",
         7.7, 2.85, 4.9, 2.5,
         size=13, color=BLUE_LIGHT, wrap=True)

add_text(slide, "สำรองข้อมูล:", 7.7, 5.1, 4.9, 0.4,
         size=12, bold=True, color=WHITE)
add_rect(slide, 7.7, 5.5, 4.9, 0.6, fill=RGBColor(0x0F, 0x17, 0x2A))
add_text(slide, "docker compose cp backend:/app/data/\\n  prompt-berk.db ./backup.db",
         7.9, 5.55, 4.6, 0.5, size=10, color=EMERALD, font="Courier New")


# ════════════════════════════════════════════════════════
# SLIDE 12 — การเข้าถึงระบบ
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=RGBColor(0xF0, 0xF7, 0xFF))
section_header(slide, "11  |  การเข้าถึงระบบ", color=BLUE_MID)

slide_title(slide, "วิธีเข้าใช้งานระบบ",
            subtitle="ผู้ใช้ทุกคนเข้าผ่านเบราว์เซอร์ — ไม่ต้องติดตั้งอะไรเพิ่ม",
            title_color=SLATE_DARK, sub_color=SLATE_MID)

access_cards = [
    ("🌐", "URL เข้าใช้งาน",
     "http://<IP-เซิร์ฟเวอร์>:8080/\n\nแชร์ลิงก์นี้ให้ผู้ใช้ทุกคนในองค์กร",
     BLUE_LIGHT, BLUE_MID),
    ("🔒", "เงื่อนไขการเข้าถึง",
     "ต้องอยู่ในเครือข่าย Intranet\nหรือเชื่อมต่อ VPN ขององค์กร",
     AMBER_LT, AMBER),
    ("👥", "ข้อมูลร่วมกัน",
     "ทุกคนเห็นรายการเบิกและ\nหมวดหมู่งบประมาณร่วมกัน",
     EMERALD_LT, EMERALD),
]

for i, (icon, title, desc, bg, accent) in enumerate(access_cards):
    lx = 0.6 + i * 4.2
    add_rect(slide, lx, 2.2, 3.8, 4.5, fill=bg)
    add_rect(slide, lx, 2.2, 3.8, 0.5, fill=accent)
    add_text(slide, icon,  lx+0.2, 2.25, 0.5, 0.45, size=18, color=WHITE)
    add_text(slide, title, lx+0.8, 2.25, 2.8, 0.45, size=14, bold=True, color=WHITE)
    add_text(slide, desc,  lx+0.3, 2.9,  3.3, 3.0,  size=14, color=SLATE_DARK, wrap=True)


# ════════════════════════════════════════════════════════
# SLIDE 13 — ข้อจำกัดและแผนถัดไป
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=SLATE_LIGHT)
section_header(slide, "12  |  ข้อจำกัดและแผนพัฒนา", color=AMBER)

slide_title(slide, "ข้อจำกัดและแผนถัดไป",
            title_color=SLATE_DARK)

# Left — limitations
add_rect(slide, 0.5, 2.0, 5.9, 5.0, fill=WHITE)
add_rect(slide, 0.5, 2.0, 5.9, 0.5, fill=AMBER)
add_text(slide, "⚠️  ข้อจำกัดปัจจุบัน", 0.75, 2.07, 5.4, 0.4,
         size=14, bold=True, color=WHITE)

limits = [
    "ไม่มีระบบ Login — ใครรู้ URL เข้าได้",
    "Longdo API Key ฝังในโค้ด\n(ควรหารือ IT หาก URL เปิดกว้าง)",
    "ไม่มีการแบ่งสิทธิ์ผู้ใช้\n(Admin = ใครก็ได้)",
    "รองรับการ deploy ที่ root path\n(ไม่รองรับ subpath เช่น /app/)",
]
add_bullet_box(slide, limits, 0.75, 2.65, 5.4, 4.0,
               size=13, color=SLATE_DARK)

# Right — roadmap
add_rect(slide, 6.9, 2.0, 5.9, 5.0, fill=WHITE)
add_rect(slide, 6.9, 2.0, 5.9, 0.5, fill=EMERALD)
add_text(slide, "🚀  แผนพัฒนาถัดไป", 7.15, 2.07, 5.4, 0.4,
         size=14, bold=True, color=WHITE)

roadmap = [
    "ระบบ Login และสิทธิ์ผู้ใช้",
    "Dashboard สรุปงบประมาณรายเดือน",
    "Export Excel นอกเหนือจาก PDF",
    "ระบบแจ้งเตือนสถานะผ่าน Email",
    "รองรับ HTTPS / SSL Certificate",
]
add_bullet_box(slide, roadmap, 7.15, 2.65, 5.4, 4.0,
               size=13, color=SLATE_DARK)


# ════════════════════════════════════════════════════════
# SLIDE 14 — สรุปและ Q&A
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=BLUE_DARK)
add_rect(slide, 0, 5.5, 13.33, 2.0, fill=BLUE_MID)

add_text(slide, "สรุป", 0.8, 0.8, 11.0, 0.8,
         size=42, bold=True, color=WHITE)

summary_points = [
    "✅  Prompt-Berk ช่วยลดขั้นตอนการเบิกจ่ายจากกระดาษมาสู่ระบบดิจิทัล",
    "✅  ข้อมูลเก็บส่วนกลาง ทุกคนเห็นร่วมกัน Real-time",
    "✅  Deploy ด้วย Docker — พร้อมใช้งานบน Intranet ทันที",
    "✅  ไม่ต้องติดตั้งโปรแกรม — ใช้ผ่านเบราว์เซอร์ได้เลย",
]
add_bullet_box(slide, summary_points, 0.8, 1.8, 11.5, 3.2,
               size=16, color=WHITE, bullet="")

add_text(slide, "❓  Questions & Answers", 0.8, 5.6, 11.5, 0.6,
         size=22, bold=True, color=WHITE)
add_text(slide, "สถาบันรับรองคุณภาพสถานพยาบาล (สรพ.)  |  Prompt-Berk  |  2026",
         0.8, 6.35, 11.5, 0.4, size=12, color=BLUE_LIGHT)


# ════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════
out_path = r"D:\Manny\promt-beark\Prompt-Berk-Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
